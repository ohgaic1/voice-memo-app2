import streamlit as st
import tempfile
import os
import re
import json
from pathlib import Path
from datetime import datetime
import shutil
import subprocess
from openai import OpenAI
import requests

# ═══════════════════════════════════════════
# 共有ライブラリ（shared-lib）を読む
# ═══════════════════════════════════════════
#
# ★2026-08-20: このアプリを Streamlit Community Cloud から PC-B へ移した。
#   外では使わないと決めたので、同梱していた分割器（_rt_all）をやめ、
#   ★他の20本以上のスクリプトと同じく shared-lib を読む形に揃えた。
#   同梱していた理由は「クラウドに shared-lib が無い」だったので、
#   ★動く場所が PC-B だけになった時点でその理由は消えている。
#
# ★try/except で「読めたら共有・読めなければ同梱」にはしない。
#   それだと環境ごとに違う実装が走り、テストが見ているものと本番で動くものが
#   別になる（2026-08-19 の事故がその形だった）。読めなければ落とす。
#
# ★探し方は「隣」を先に見る。C:\dev\voice-memo-app2 の隣が C:\dev\shared-lib
#   なので、PC-B でも CI でも同じ形（リポジトリを並べて置く）で解決できる。
#
# ★破綻点（先に挙げたもの）:
#   隣に shared-lib が無い環境では、ここで ModuleNotFoundError になり起動しない。
#   今日の事故（クラウドに shared-lib が無い）の★逆方向。
#   検知: .github/workflows/ci-cd.yml が ubuntu で shared-lib を隣に checkout し、
#   ★実際に起動して health 200 を確かめる。checkout の行を外すとそこが赤になる。
#   PC-B 側は dev-tools/tests/test_voice_memo_shared_lib.py が毎日確かめる。
import sys as _sys

for _sl in (Path(__file__).resolve().parent.parent / "shared-lib",
            Path(r"C:\dev\shared-lib"),
            Path(r"C:\Users\ohgai\dev\shared-lib")):
    if _sl.is_dir():
        if str(_sl) not in _sys.path:
            _sys.path.insert(0, str(_sl))
        break

#: ★Notion の rich_text は1要素 2000字まで。`text[:2000]` は保存時に中身を捨てる
#  （切り捨てた事実は保存値に現れない）。分割は shared-lib に1つだけ置く。
#  ★ここで自前で書かないこと。切り捨てにも戻さないこと。
#  旧実装は 1990字ごとに割って先頭10塊で打ち切り、呼び出し側も text[:2000] で
#  先に切っていた（19,900字で黙って消えていた）。
from notion_rich_text import text_to_rich_text as _rt_all   # noqa: E402

# ★.env の読み込みは shared-lib/env_loader を使う（override=True）。
#   素の load_dotenv() は使わない ── 古い永続環境変数が .env に勝ち、
#   2026-07-23〜27 に NOTION_TOKEN が4日間 401 を出し続けた事故の原因。
#   ★これで OpenAI / Notion のキーを画面で毎回入力する必要がなくなる。
from env_loader import load_env                             # noqa: E402
import api_usage                                            # noqa: E402
import api_prices                                           # noqa: E402

# ★このリポジトリの中の純粋な部分（API を1回も呼ばない）。試験はこちらで行う。
import report_builder as RB                                 # noqa: E402

#: 使用量の記録の「どこから」に入るアプリ名。
USAGE_APP = "voice-memo-app2"

# ★音声認識に渡す分野語彙。vocab/legal_ja.txt から読む（コードに書かない）。
#   .env の VOCAB_FILE で別のファイルを足せる。
#   ★入りきらなかった数も持っておき、画面に出す（黙って捨てない）。
VOCAB_TERMS = RB.load_vocab(Path(__file__).resolve().parent,
                            extra_path=os.environ.get("VOCAB_FILE", ""))
VOCAB_PROMPT, VOCAB_USED, VOCAB_DROPPED = RB.whisper_prompt(VOCAB_TERMS)

# ★レポート生成に渡す誤変換の対応表（音声認識用の語彙とは役割が違う）。
#   2026-08-21 実測: 文字起こし段では語彙を渡しても 公権45回 が直らず、
#   レポート段で 138→18 まで減った。残りを減らすための材料。
#   ★機械的に置換はしない。文脈を見てモデルが直す。
CORRECTIONS = RB.load_corrections(Path(__file__).resolve().parent,
                                  extra_path=os.environ.get("CORRECTIONS_FILE", ""))

load_env()

# ─────────────────────────────────────────
# ページ設定
# ─────────────────────────────────────────
st.set_page_config(
    page_title="音声メモアプリ Pro",
    page_icon="🎙️",
    layout="wide"
)

if "api_key" not in st.session_state:
    # 優先順位: 1. 環境変数 / 2. st.secrets / 3. 空文字（サイドバー入力へ）
    _key = os.environ.get("OPENAI_API_KEY", "")
    if not _key:
        try:
            _key = st.secrets.get("OPENAI_API_KEY", "")
        except Exception:
            _key = ""
    st.session_state.api_key = _key
if "results" not in st.session_state:
    st.session_state.results = []

# Notion設定（環境変数から読み込み）
# ★2026-08-21: 名前を NOTION_TOKEN に揃えた。
#   それまで NOTION_API_KEY だけを探しており、共有の .env に在るのは
#   ★NOTION_TOKEN なので、画面に「NOTION_API_KEY 未設定」と出て
#   研修DBに登録できなかった。
#   実測: .env の Notion 系のキー名は NOTION_TOKEN と各DBのIDのみ。
#         NOTION_API_KEY という名前は★存在しない。
#   他のシステム（voice_report_batch / kenshu_to_dify_layer3 / bot）は
#   すべて NOTION_TOKEN を使っている。★このアプリだけが違っていた。
#   ★古い名前も残して読む（どちらでも動く。移行中に壊さないため）。
NOTION_API_KEY     = (os.environ.get("NOTION_TOKEN", "")
                      or os.environ.get("NOTION_API_KEY", ""))
NOTION_DB_KENSHU   = os.environ.get("NOTION_DB_ID_KENSHU", "475162c3cf1f4993a9b231e202ec40fb")


# ═══════════════════════════════════════════
# トークン制限対策：長いテキストを事前圧縮
# ═══════════════════════════════════════════
# ★2026-08-20: MAX_TRANSCRIPT_CHARS（12000字で打ち切り）を廃止した。
#   一定の長さを超えると先頭・中間・末尾の3か所を抜き取るだけで、
#   残りを一度も読まなかった。今日の2時間の講演では★77.4%が読まれず、
#   欠落を告げる仕組みが無いので出力は正常に見えていた。
#   ★いまは区画に分けて全文を読み、読んだ量と捨てた量をレポートに書く。
#   区画の大きさは report_builder.CHUNK_CHARS。
MAX_MATERIAL_CHARS   = 12000   # 資料テキストの上限（★超えた分は画面に出す）


# ★compress_transcript（先頭・中間・末尾を4000字ずつ抜き取って要約する関数）は
#   2026-08-20 に削除した。★これが 77.4% を黙って捨てていた実体。
#   代わりに generate_report_full() が全文を区画に分けて読む。
#   ★戻さないこと。戻したら tests/test_reads_everything.py が落ちる。


# ═══════════════════════════════════════════
# 音声処理
# ═══════════════════════════════════════════
class FfmpegNotFound(RuntimeError):
    """★ffmpeg が見つからない。黙って進まず、ここで止める。"""


def _ffmpeg() -> str:
    """★ffmpeg の実体を1か所で解決する。

    ★2026-08-20 実測: PC-B の PATH に ffmpeg は無い。
      Streamlit Community Cloud では packages.txt（中身は ffmpeg の1行）が
      入れてくれていたので、★PC-B へ移すとここが失われる。
      移設で消える依存はこれだけ。

    探す順: .env の FFMPEG_PATH → PATH。
    ★見つからないときは False を返さず例外で止める。
      戻り値で失敗を伝えると呼び出し側が「圧縮エラー」としか出せず、
      何が足りないのか画面から分からない。
    """
    p = (os.environ.get("FFMPEG_PATH", "") or "").strip().strip('"')
    if p and Path(p).is_file():
        return p
    found = shutil.which("ffmpeg")
    if found:
        return found
    raise FfmpegNotFound(
        "ffmpeg が見つかりません。24MB を超える音声の圧縮・分割ができません。 "
        "ffmpeg を入れて PATH を通すか、"
        ".env に FFMPEG_PATH=<ffmpeg.exe のフルパス> を書いてください。")


def _ffprobe() -> str:
    """★ffprobe も同じ理由で PC-B の PATH に無い（2026-08-20 実測）。

    ffmpeg の公式ビルドには同じ場所に入っているので、まず ffmpeg の隣を見る。
    """
    p = (os.environ.get("FFPROBE_PATH", "") or "").strip().strip('"')
    if p and Path(p).is_file():
        return p
    found = shutil.which("ffprobe")
    if found:
        return found
    try:
        near = Path(_ffmpeg()).with_name("ffprobe.exe")
        if near.is_file():
            return str(near)
    except FfmpegNotFound:
        pass
    raise FfmpegNotFound(
        "ffprobe が見つかりません。音声の分割ができません。 "
        "ffmpeg を入れると同じ場所に入ります（PATH を通すか "
        ".env に FFPROBE_PATH を書いてください）。")


def compress_audio(input_path, output_path):
    try:
        subprocess.run(
            [_ffmpeg(), "-i", input_path, "-vn", "-ac", "1",
             "-ar", "16000", "-b:a", "32k", "-y", output_path],
            check=True, capture_output=True
        )
        return True
    except Exception as e:
        st.error(f"圧縮エラー: {e}")
        return False


def split_audio(input_path, chunk_sec=600):
    """音声を10分チャンクに分割"""
    chunks = []
    try:
        # 時間を取得
        probe = subprocess.run(
            [_ffprobe(), "-v", "error", "-show_entries", "format=duration",
             "-of", "default=noprint_wrappers=1:nokey=1", input_path],
            capture_output=True, text=True, check=True
        )
        duration = float(probe.stdout.strip())
        num_chunks = int(duration / chunk_sec) + 1
        
        # 拡張子に関わらず.mp3で統一
        base = os.path.splitext(input_path)[0]
        
        for i in range(num_chunks):
            output_chunk = f"{base}_chunk{i}.mp3"
            subprocess.run(
                [_ffmpeg(), "-i", input_path,
                 "-ss", str(i * chunk_sec), "-t", str(chunk_sec),
                 "-c", "copy", "-y", output_chunk],
                check=True, capture_output=True
            )
            if os.path.exists(output_chunk) and os.path.getsize(output_chunk) > 1000:
                chunks.append(output_chunk)
        
        return chunks
    except Exception as e:
        st.error(f"分割エラー: {e}")
        return []


def _audio_minutes(path) -> float:
    """★音声の長さ（分）。取れなければ 0 を返す（推測しない）。"""
    try:
        r = subprocess.run(
            [_ffprobe(), "-v", "error", "-show_entries", "format=duration",
             "-of", "default=noprint_wrappers=1:nokey=1", str(path)],
            capture_output=True, text=True, check=True)
        return float(r.stdout.strip()) / 60.0
    except Exception:                                        # noqa: BLE001
        return 0.0


def _segments_of(resp) -> list:
    """★verbose_json の segments を辞書の一覧で取り出す。

    SDK は pydantic の物を返すので、そのままでは添字で読めない。
    ★取れなければ空を返す。ここで時刻を作らない（無い物を作らない）。
    """
    try:
        d = resp.model_dump()
    except Exception:                                        # noqa: BLE001
        d = resp if isinstance(resp, dict) else {}
    segs = d.get("segments") or []
    return [x for x in segs if isinstance(x, dict)]


def transcribe_audio(file_path, api_key, vocab_prompt: str = ""):
    """音声を文字起こしする。返す: (時刻付き本文, セグメント, 経過の記録)。

    ★2026-08-20 追加。どちらも★渡すだけ・受け取るだけで、費用は変わらない。
      ・vocab_prompt … 分野語彙のヒント。whisper が受け付けるのに渡していなかった。
      ・response_format="verbose_json" … 時刻付きのセグメントが返る。
        これで章や欠落の位置を★時刻で言えるようになる。

    ★2026-08-21 変更。それまでは「まるごと圧縮 → それでも大きければ分割」だった。
      2時間の音声では 69kbps → 32kbps へ★全体を再エンコードしてから分割しており、
      分割すれば1区画あたり上限を下回るのに、無駄に音質を落としていた。
      ★いまは「先に分割 → 上限を超える区画だけ圧縮」にする。
      圧縮した区画があれば、それが何本かを★記録してレポートに書く
      （黙って音質を落とさないため）。

    ★2026-08-21 の実測メモ: 圧縮の有無で認識の質に有意な差は見られなかった
      （誤変換語の合計 166 vs 190・方向はむしろ圧縮側が良い／スペクトルの
       95%帯域 1770Hz vs 1777Hz）。それでも★不要な再エンコードはしない。
      得るものが無く、失うもの（時間と音質）があるため。
    """
    client = OpenAI(api_key=api_key)
    max_size = 24 * 1024 * 1024
    info = {"chunks": 0, "compressed": [], "note": ""}

    try:
        size = os.path.getsize(file_path)

        # ── 上限内ならそのまま1回で ──
        if size <= max_size:
            info["chunks"] = 1
            with open(file_path, "rb") as f:
                with api_usage.record(app=USAGE_APP, site="transcribe.single",
                                      provider="openai",
                                      operation="audio.transcriptions",
                                      model="whisper-1") as _rec:
                    resp = client.audio.transcriptions.create(
                        model="whisper-1", file=f, language="ja",
                        prompt=vocab_prompt or "",
                        response_format="verbose_json")
                    _rec.sdk_response(resp)
            segs = _segments_of(resp)
            return ((RB.timestamped_text(segs) or getattr(resp, "text", "")),
                    segs, info)

        # ── ★先に分割する（原音のまま）──
        st.info("  ✂️ 分割中...")
        chunks = split_audio(file_path)
        if not chunks:
            return None, [], info
        info["chunks"] = len(chunks)

        texts = []
        all_segments = []
        pb = st.progress(0)
        for i, chunk in enumerate(chunks):
            work = chunk
            tmp_comp = None
            # ★上限を超える区画だけ圧縮する。超えていなければ触らない。
            if os.path.getsize(chunk) > max_size:
                tmp_comp = os.path.splitext(chunk)[0] + "_comp.mp3"
                if not compress_audio(chunk, tmp_comp):
                    st.error("  ❌ 区画%d の圧縮に失敗しました" % i)
                    os.remove(chunk)
                    continue
                work = tmp_comp
                info["compressed"].append(i)

            with open(work, "rb") as f:
                with api_usage.record(app=USAGE_APP, site="transcribe.chunk",
                                      provider="openai",
                                      operation="audio.transcriptions",
                                      model="whisper-1") as _rec:
                    resp = client.audio.transcriptions.create(
                        model="whisper-1", file=f, language="ja",
                        prompt=vocab_prompt or "",
                        response_format="verbose_json")
                    _rec.sdk_response(resp)
            segs = _segments_of(resp)
            off = i * 600            # split_audio の1区画は600秒
            all_segments.extend(
                dict(sg, start=float(sg.get("start", 0)) + off,
                     end=float(sg.get("end", 0)) + off) for sg in segs)
            texts.append(RB.timestamped_text(segs, offset_sec=off)
                         or getattr(resp, "text", ""))
            pb.progress((i + 1) / len(chunks))
            os.remove(chunk)
            if tmp_comp and os.path.exists(tmp_comp):
                os.remove(tmp_comp)

        info["note"] = RB.compression_note(info)
        if info["compressed"]:
            st.warning("⚠️ " + info["note"])
        return "\n".join(texts).strip(), all_segments, info

    except Exception as e:
        st.error(f"文字起こしエラー: {e}")
        return None, [], info


# ═══════════════════════════════════════════
# 資料テキスト抽出
# ═══════════════════════════════════════════
def extract_pdf_text(file_path):
    try:
        import pdfplumber
        texts = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    texts.append(t)
        return "\n".join(texts)
    except Exception:
        try:
            from pypdf import PdfReader
            reader = PdfReader(file_path)
            return "\n".join(p.extract_text() or "" for p in reader.pages)
        except Exception as e:
            return f"[PDF読み取りエラー: {e}]"


def extract_pptx_text(file_path):
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        texts = []
        for i, slide in enumerate(prs.slides, 1):
            parts = [s.text.strip() for s in slide.shapes
                     if hasattr(s, "text") and s.text.strip()]
            if parts:
                texts.append(f"【スライド{i}】\n" + "\n".join(parts))
        return "\n\n".join(texts)
    except Exception as e:
        return f"[PPTX読み取りエラー: {e}]"


def extract_docx_text(file_path):
    try:
        from docx import Document
        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        return f"[DOCX読み取りエラー: {e}]"


def extract_material_text(uploaded_file):
    suffix = Path(uploaded_file.name).suffix.lower()
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as f:
        f.write(uploaded_file.read())
        tmp = f.name
    try:
        if suffix == ".pdf":
            return extract_pdf_text(tmp)
        elif suffix in [".pptx", ".ppt"]:
            return extract_pptx_text(tmp)
        elif suffix in [".docx", ".doc"]:
            return extract_docx_text(tmp)
        return f"[未対応形式: {suffix}]"
    finally:
        if os.path.exists(tmp):
            os.remove(tmp)


# ═══════════════════════════════════════════
# YouTube字幕取得
# ═══════════════════════════════════════════
def extract_youtube_id(url: str) -> str | None:
    patterns = [
        r'(?:v=|/v/|youtu\.be/|/embed/)([^&\n?#]+)',
        r'(?:shorts/)([^&\n?#]+)',
    ]
    for pattern in patterns:
        m = re.search(pattern, url)
        if m:
            return m.group(1)
    return None


def _fetch_entries_text(entries) -> str:
    """FetchedTranscript エントリ（v1.x オブジェクト or v0.x 辞書）からテキストを結合"""
    parts = []
    for e in entries:
        if hasattr(e, "text"):
            parts.append(e.text)
        elif isinstance(e, dict):
            parts.append(e.get("text", ""))
    return " ".join(parts)


def _youtube_whisper_fallback(url: str, video_id: str) -> tuple:
    """字幕不可の場合に yt-dlp で音声ダウンロード → Whisper 文字起こしを試みる"""
    try:
        import yt_dlp  # optional dependency
    except ImportError:
        return (
            None,
            f"字幕が見つかりません（ID: {video_id}）。\n"
            "音声からの文字起こしを行うには yt-dlp が必要です: pip install yt-dlp",
        )

    st.info("💬 字幕が見つかりません。yt-dlp で音声をダウンロードして Whisper で文字起こしを試みます...")
    try:
        import tempfile, glob
        with tempfile.TemporaryDirectory() as tmpdir:
            out_tmpl = os.path.join(tmpdir, "%(id)s.%(ext)s")
            ydl_opts = {
                "format": "bestaudio[ext=m4a]/bestaudio/best",
                "outtmpl": out_tmpl,
                "quiet": True,
                "no_warnings": True,
            }
            with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                ydl.download([url])
            files = glob.glob(os.path.join(tmpdir, "*"))
            if not files:
                return None, "yt-dlp でのダウンロードに失敗しました"
            audio_path = files[0]
            api_key = st.session_state.get("api_key", "")
            if not api_key:
                return None, "音声ダウンロード成功しましたが OpenAI API キーがないため Whisper 文字起こしができません"
            transcript, _segs, _info = transcribe_audio(audio_path, api_key,
                                                        VOCAB_PROMPT)
            if transcript:
                st.success(f"✅ Whisper フォールバック成功（{len(transcript):,}文字）")
                return transcript, video_id
            return None, "Whisper 文字起こしに失敗しました"
    except Exception as e:
        return None, f"音声ダウンロード / 文字起こしエラー: {str(e)[:200]}"


def get_youtube_transcript(url: str) -> tuple:
    """YouTube URL から字幕テキストを取得。字幕不可なら Whisper フォールバック。
    Returns: (text_or_None, video_id_or_error_message)
    """
    # ── パッケージ確認 ──
    try:
        from youtube_transcript_api import YouTubeTranscriptApi
        from youtube_transcript_api._errors import (
            NoTranscriptFound,
            TranscriptsDisabled,
            VideoUnavailable,
            InvalidVideoId,
        )
    except ImportError as ie:
        return None, f"youtube-transcript-api がインストールされていません: {ie}"

    # v1.x の新エラー（v0.x では存在しない場合がある）
    try:
        from youtube_transcript_api._errors import RequestBlocked, IpBlocked, AgeRestricted
        _HAS_V1_ERRORS = True
    except ImportError:
        RequestBlocked = IpBlocked = AgeRestricted = Exception  # 型チェック用ダミー
        _HAS_V1_ERRORS = False

    # ── URL からビデオ ID 抽出 ──
    video_id = extract_youtube_id(url)
    if not video_id:
        return None, (
            "YouTube URL からビデオ ID を抽出できませんでした。\n"
            "対応形式: https://www.youtube.com/watch?v=XXXXX "
            "/ https://youtu.be/XXXXX / Shorts URL"
        )

    # ── 字幕取得（v1.x インスタンスメソッド） ──
    try:
        api = YouTubeTranscriptApi()
        transcript_list = api.list(video_id)

        # 優先度: ①手動日本語 → ②自動生成日本語 → ③英語 → ④何でも最初の1件
        transcript = None
        for finder, langs in [
            ("find_manually_created_transcript", ["ja", "ja-JP"]),
            ("find_generated_transcript",        ["ja", "ja-JP"]),
            ("find_transcript",                  ["en", "en-US", "en-GB"]),
        ]:
            try:
                transcript = getattr(transcript_list, finder)(langs)
                break
            except Exception:
                continue

        if transcript is None:
            for t in transcript_list:
                transcript = t
                break

        if transcript is None:
            return _youtube_whisper_fallback(url, video_id)

        entries = transcript.fetch()
        text = _fetch_entries_text(entries)
        lang = getattr(transcript, "language", "不明")
        is_gen = getattr(transcript, "is_generated", False)
        st.caption(f"📝 字幕言語: {lang}（{'自動生成' if is_gen else '手動作成'}）")
        return text, video_id

    except TranscriptsDisabled:
        st.warning("⚠️ この動画では字幕が無効です。Whisper フォールバックを試みます...")
        return _youtube_whisper_fallback(url, video_id)

    except NoTranscriptFound:
        st.warning("⚠️ 字幕が見つかりません。Whisper フォールバックを試みます...")
        return _youtube_whisper_fallback(url, video_id)

    except VideoUnavailable:
        return None, f"動画が存在しないか非公開です（ID: {video_id}）"

    except InvalidVideoId:
        return None, f"無効なビデオ ID です: {video_id}"

    except Exception as e:
        err_str = str(e)
        # v1.x 特有のエラーを文字列でもチェック（_HAS_V1_ERRORS が False の場合の保険）
        if _HAS_V1_ERRORS and isinstance(e, AgeRestricted):
            return None, "年齢制限のある動画は字幕取得できません（認証未対応）"
        if _HAS_V1_ERRORS and isinstance(e, (RequestBlocked, IpBlocked)):
            return None, (
                "YouTube にアクセスがブロックされました（IP 制限）。\n"
                "しばらく待つか、別のネットワークから試してください。"
            )
        if "429" in err_str or "too many requests" in err_str.lower():
            return None, "リクエスト過多です（429）。しばらく待ってから再試行してください。"
        if "age" in err_str.lower() and "restrict" in err_str.lower():
            return None, "年齢制限のある動画は字幕取得できません"
        if "unavailable" in err_str.lower() or "private" in err_str.lower():
            return None, f"動画が利用不可または非公開です: {err_str[:100]}"
        return None, f"字幕取得エラー: {err_str[:200]}"


# ═══════════════════════════════════════════
# レポートからタイトル・タグを抽出
# ═══════════════════════════════════════════
def extract_title_from_report(report: str) -> str:
    for line in report.splitlines():
        line = line.strip()
        if line.startswith("# ") and not line.startswith("## "):
            return line[2:].strip()
    return "音声メモレポート"


def extract_tags_from_report(report: str) -> list[str]:
    for line in report.splitlines():
        line = line.strip()
        if line.startswith("> タグ：") or line.startswith("> タグ:"):
            tags_str = re.sub(r'^> タグ[：:]', '', line).strip()
            return [t.strip() for t in re.split(r'[、,，　 ]+', tags_str) if t.strip()]
    return []


# ═══════════════════════════════════════════
# Markmap生成
# ═══════════════════════════════════════════
def generate_markmap(report: str, api_key: str) -> str | None:
    """PLAUDレポートからMarkmap用Markdown見出し構造を生成"""
    client = OpenAI(api_key=api_key)
    try:
        resp = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "あなたはMarkdown見出し構造のみを出力する専門家です。コードブロックや説明文は一切不要です。"},
                {"role": "user", "content": f"""以下のレポートをMarkmap形式のMarkdown見出し構造に変換してください。
ルール：
- 見出し（# ## ###）のみ使用
- 各ノードは短いキーワード（15文字以内）
- コードブロック不要、見出しのみ出力
- 深さは最大3階層

レポート（抜粋）:
{report[:5000]}"""}
            ],
            temperature=0.3,
            max_tokens=800,
        )
        return resp.choices[0].message.content
    except Exception as e:
        return None


def render_markmap_html(markmap_md: str) -> str:
    return f"""<!DOCTYPE html>
<html>
<head>
<meta charset="UTF-8">
<script src="https://cdn.jsdelivr.net/npm/markmap-autoloader@latest"></script>
<style>
html,body{{margin:0;padding:0;width:100%;height:600px;overflow:hidden;}}
#mm{{width:100%;height:600px;}}
</style>
</head>
<body>
<div class="markmap" id="mm">

{markmap_md}

</div>
<script>
window.addEventListener('load', function() {{
  setTimeout(function() {{
    var svgs = document.querySelectorAll('svg.markmap');
    if (svgs.length > 0) {{
      try {{
        var mmInstance = svgs[0].__markmap;
        if (mmInstance && mmInstance.fit) {{ mmInstance.fit(); }}
      }} catch(e) {{}}
    }}
  }}, 800);
}});
</script>
</body>
</html>"""


# ═══════════════════════════════════════════
# Notion ブロックヘルパー
# ═══════════════════════════════════════════
def _rich_text(content: str) -> list:
    """★分割は shared-lib/notion_rich_text に1つだけ置く。ここで書かない。

    旧実装は 1990字ごとに割ったうえで★先頭10塊で打ち切っていた（19,900字で
    黙って消える）。さらに呼び出し側が `text[:2000]` で先に切っていたため、
    ★2000字を超える段落は保存すらされていなかった。
    分割器は 2000字ごとに割り、入らない長さは例外で★はっきり失敗させる。

    ★2026-08-20: 一度は同梱に戻した（Streamlit Community Cloud に shared-lib が
      無く起動しなかったため）。同日 PC-B へ移したので★同梱をやめ、共有に戻した。
      経緯は冒頭のコメント。
    """
    return _rt_all(content)

def _heading_block(level: int, text: str) -> dict:
    t = f"heading_{level}"
    return {"object": "block", "type": t, t: {"rich_text": _rich_text(text[:200])}}

def _paragraph_block(text: str) -> dict:
    return {"object": "block", "type": "paragraph",
            "paragraph": {"rich_text": _rich_text(text)}}

def _bulleted_block(text: str) -> dict:
    return {"object": "block", "type": "bulleted_list_item",
            "bulleted_list_item": {"rich_text": _rich_text(text)}}

def _numbered_block(text: str) -> dict:
    return {"object": "block", "type": "numbered_list_item",
            "numbered_list_item": {"rich_text": _rich_text(text)}}

def _quote_block(text: str) -> dict:
    return {"object": "block", "type": "quote",
            "quote": {"rich_text": _rich_text(text)}}

def _divider_block() -> dict:
    return {"object": "block", "type": "divider", "divider": {}}

def _code_block(content: str, language: str = "markdown") -> dict:
    return {"object": "block", "type": "code",
            "code": {"rich_text": _rich_text(content), "language": language}}


def markdown_to_notion_blocks(md: str) -> list:
    """マークダウン文字列をNotionブロックリストに変換"""
    blocks = []
    for line in md.splitlines():
        if line.startswith("# "):
            blocks.append(_heading_block(1, line[2:].strip()))
        elif line.startswith("## "):
            blocks.append(_heading_block(2, line[3:].strip()))
        elif line.startswith("### "):
            blocks.append(_heading_block(3, line[4:].strip()))
        elif re.match(r'^\d+\. ', line):
            text = re.sub(r'^\d+\. ', '', line).strip()
            blocks.append(_numbered_block(text))
        elif line.startswith("- ") or line.startswith("* "):
            text = line[2:].strip()
            blocks.append(_bulleted_block(text))
        elif line.startswith("> "):
            text = line[2:].strip()
            blocks.append(_quote_block(text))
        elif line.strip() == "---":
            blocks.append(_divider_block())
        elif line.strip():
            text = re.sub(r'`([^`]+)`', r'\1', line)
            text = re.sub(r'\*\*([^*]+)\*\*', r'\1', text)
            text = re.sub(r'\*([^*]+)\*', r'\1', text)
            blocks.append(_paragraph_block(text))
    return blocks


def _append_blocks(page_id: str, blocks: list, headers: dict) -> None:
    """90ブロックのバッチに分けてNotionページに追記"""
    for i in range(0, len(blocks), 90):
        batch = blocks[i:i+90]
        resp = requests.patch(
            f"https://api.notion.com/v1/blocks/{page_id}/children",
            headers=headers,
            json={"children": batch},
            timeout=30,
        )
        if not resp.ok:
            raise ValueError(f"ブロック追記エラー {resp.status_code}: {resp.text[:200]}")


# ═══════════════════════════════════════════
# Notion 研修DB への保存
# ═══════════════════════════════════════════
def save_to_notion_kenshu(
    title: str,
    tags: list,
    source_type: str,
    report: str,
    summary: str,
    transcript: str = "",
    markmap_md: str = "",
    summary_data: dict = None,
    source_info: list = None,
    attachment_file_info: list = None,
) -> bool:
    if not NOTION_API_KEY:
        st.error("⚠️ NOTION_TOKEN が未設定です。OneDrive の secrets フォルダにある .env を確認してください。")
        return False

    headers = {
        "Authorization": f"Bearer {NOTION_API_KEY}",
        "Notion-Version": "2022-06-28",
        "Content-Type": "application/json",
    }
    date_iso = datetime.now().strftime("%Y-%m-%d")

    properties: dict = {
        "タイトル": {"title": [{"text": {"content": title[:200]}}]},
        "ジャンル": {"select": {"name": "その他"}},
        "種別":   {"select": {"name": source_type}},
        "実施日": {"date": {"start": date_iso}},
        "作成日": {"date": {"start": date_iso}},
        "概要":   {"rich_text": [{"text": {"content": summary[:500]}}]},
    }
    if tags:
        properties["タグ"] = {"multi_select": [{"name": t[:100]} for t in tags[:5]]}

    # ── ページ作成（本文なし）──
    try:
        resp = requests.post(
            "https://api.notion.com/v1/pages",
            headers=headers,
            json={"parent": {"database_id": NOTION_DB_KENSHU}, "properties": properties},
            timeout=30,
        )
        if not resp.ok:
            msg = resp.json().get("message", resp.text[:200])
            st.error(f"Notion保存エラー: {resp.status_code} - {msg}")
            return False
        page_data = resp.json()
        page_id = page_data["id"]
        page_url = page_data.get("url", "")
    except Exception as e:
        st.error(f"Notion保存エラー: {e}")
        return False

    try:
        # ① 元データ
        src_blocks = [_heading_block(2, "① 元データ")]
        if source_info:
            for s in source_info:
                src_blocks.append(_bulleted_block(s))
        else:
            src_blocks.append(_paragraph_block("（情報なし）"))
        _append_blocks(page_id, src_blocks, headers)

        # ② 添付資料
        att_blocks = [_divider_block(), _heading_block(2, "② 添付資料")]
        if attachment_file_info:
            for fi in attachment_file_info:
                size_kb = fi.get("size", 0) // 1024
                att_blocks.append(_bulleted_block(f"{fi['name']}  ({size_kb} KB)"))
        else:
            att_blocks.append(_paragraph_block("（なし）"))
        _append_blocks(page_id, att_blocks, headers)

        # ③ 文字起こし（トグルブロック）
        toggle_header = [
            _divider_block(),
            {"object": "block", "type": "toggle",
             "toggle": {"rich_text": _rich_text("③ 文字起こし（クリックで展開）")}},
        ]
        tog_resp = requests.patch(
            f"https://api.notion.com/v1/blocks/{page_id}/children",
            headers=headers,
            json={"children": toggle_header},
            timeout=30,
        )
        if tog_resp.ok and transcript.strip():
            toggle_id = None
            for b in tog_resp.json().get("results", []):
                if b.get("type") == "toggle":
                    toggle_id = b["id"]
                    break
            if toggle_id:
                tr_chunks = [transcript[i:i+1990] for i in range(0, min(len(transcript), 60000), 1990)]
                tr_blocks = [_paragraph_block(c) for c in tr_chunks[:90]]
                _append_blocks(toggle_id, tr_blocks, headers)

        # ④ PLAUDレポート
        report_header = [_divider_block(), _heading_block(2, "④ PLAUDレポート")]
        _append_blocks(page_id, report_header, headers)
        report_blocks = markdown_to_notion_blocks(report)
        _append_blocks(page_id, report_blocks[:200], headers)

        # ⑤ マインドマップ
        if markmap_md:
            mm_blocks = [_divider_block(), _heading_block(2, "⑤ マインドマップ（Markmap）")]
            mm_chunks = [markmap_md[i:i+1990] for i in range(0, min(len(markmap_md), 10000), 1990)]
            for chunk in mm_chunks[:5]:
                mm_blocks.append(_code_block(chunk, "markdown"))
            _append_blocks(page_id, mm_blocks, headers)

        # ⑥ 構造化サマリー
        if summary_data:
            sum_blocks = [_divider_block(), _heading_block(2, "⑥ 構造化サマリー")]
            if summary_data.get("one_line"):
                sum_blocks.append(_paragraph_block(f"📌 {summary_data['one_line']}"))
            if summary_data.get("flow"):
                sum_blocks.append(_heading_block(3, "フロー"))
                for item in summary_data["flow"]:
                    sum_blocks.append(_bulleted_block(
                        f"【{item.get('time','')}】{item.get('topic','')} — {item.get('summary','')}"
                    ))
            if summary_data.get("decisions"):
                sum_blocks.append(_heading_block(3, "決定事項"))
                for d in summary_data["decisions"]:
                    sum_blocks.append(_bulleted_block(f"✅ {d.get('title','')}：{d.get('detail','')}"))
            if summary_data.get("actions"):
                sum_blocks.append(_heading_block(3, "アクションアイテム"))
                for a in summary_data["actions"]:
                    sum_blocks.append(_bulleted_block(
                        f"[{a.get('priority','')}] {a.get('what','')} — {a.get('who','未定')} / {a.get('when','期限未定')}"
                    ))
            if summary_data.get("keywords"):
                sum_blocks.append(_heading_block(3, "キーワード"))
                sum_blocks.append(_paragraph_block("  ".join(summary_data["keywords"])))
            _append_blocks(page_id, sum_blocks, headers)

        st.success(f"✅ Notionに保存しました！ [ページを開く]({page_url})")
        return True

    except Exception as e:
        st.error(f"Notion本文追記エラー: {e}")
        return False


# ═══════════════════════════════════════════
# GPT：Plaud風レポート
# ═══════════════════════════════════════════
# ★generate_report（文字起こしを 12000字で切ってから1回で書かせていた関数）は
#   2026-08-20 に削除した。出力も max_tokens=4000 で頭打ちだった。
#   代わりは generate_report_full()。★戻さないこと。


# ═══════════════════════════════════════════
# GPT：全文を読んでレポートを作る（★抜き取らない）
# ═══════════════════════════════════════════
_CHUNK_SYSTEM = (
    "あなたは行政書士事務所の記録係です。渡された文字起こしの区間を、"
    "★要約せずに章立てして書き起こします。"
    "★渡された範囲に書かれていることだけを使い、知識で補いません。"
    "★分からないものは「不明」「言及なし」と書き、空欄にしません。"
)

# ★章数の目安は書かない（2026-08-20 に撤回）。
#   実測で2時間の講演に24章が妥当と分かっている。数を目安に合わせるために
#   主題を混ぜると、知識ベースで章が引けなくなる。★数は結果に任せる。
_CHUNK_TEMPLATE = """以下は、ある講演の文字起こしの★一部（区間 {i}/{n}）です。
この区間で語られたことを、論点ごとに章に分けて書いてください。

【この区間の直前（文脈のためだけに渡します。★ここから章を作らないこと）】
{context}

【この区間の本文（★ここを最初から最後まで全部読んで章にする）】
{body}
{material}
---
■ 各章の形（★この順で、見出しも含めてそのまま使うこと）

### 【主題】
主題は★名詞。この章が何の話かを表す。★他の章を参照しない
（「前述の」「上記の」「先ほどの」を使わない）。

**何が変わるのか／何が論点か**
★要約しない。話の筋・理由・背景を落とさずに書く。

**講師が述べた根拠・理由**
なぜそうなるのか。制度趣旨、立法の背景、実務上の必要性など。

**実務でどう効くか**
事務所の業務（相続・後見・許認可・医療介護）にどう関わるか。
★語られていなければ「言及なし」と書く。推測で埋めない。

**講師が示した具体例**
例示があればそのまま。無ければ「なし」。

**言い切られていないこと**
「〜と思われる」「今後の議論」など、★講師がどれだけ確信しているか（確信度）。
★ここを落とさない。断定に変えない。
★下の「公表状態」とは別物。確信して話しているが未公表、ということがある。

■ ★守ること

1. ★この文章は音声認識の出力で、法律用語が高い確率で誤変換されている。
   配付資料と突き合わせて正しい法令用語に直すこと。
{corrections}
   直したものは下の「用語の訂正」に「直す前 → 直した後」で必ず記録する。
2. ★どう直せばよいか分からないものは直さず、
   【聞き取り不確か：〜】と印を付けて残す。★自然な言葉に置き換えない。
3. ★次のものは本文に書かない。下の一覧にだけ書く。
   ・日付（令和○年○月、○月○日 など）
   ・条文番号（○条、○条の○ など）
   ・金額・件数・割合
   本文では★「（→日付・数値の一覧）」と書いて参照する。
   ★ただし次は「数値」ではなく用語の一部なので、そのまま本文に書いてよい:
     第2種社会福祉事業／9条列挙行為／3類型／2本柱／2段構え／各号行為
   ★2026-08-21 実測: この指示を守れず、本文に日付2行・条文番号5行が残った。
     残っていれば機械的に数えてレポートに出る。
4. ★本文に [HH:MM:SS] の時刻を書き写さない。
   本文の行頭に付いているのは位置を知るための印であって、本文の一部ではない。
   時刻を書いてよいのは【聞き取り不確か：〜】の印を付けるときだけ。
   ★2026-08-21 実測: 時刻をそのまま写した行が14行あった。
5. ★聞き取れなかった箇所には【聞き取り不確か：〜】の印を付ける。
   ここでだけ時刻を書いてよい。
6. ★一つの章に複数の主題を混ぜない。主題が変われば章を分ける。
   ★章の数を気にしない。分けるべきなら分ける。
7. ★この区間の最後まで扱う。途中で「以下略」としない。
8. 前置きや締めの文は書かない。

■ 章の後ろに、次の3つを必ず付けること（★該当が無ければ「なし」の1行）

<<<用語の訂正>>>
直す前 → 直した後
（この区間で直したものを1行ずつ）

<<<日付・数値>>>
事項 | 値 | 明確 or 推定 or 不確か
（この区間に出た日付・数値・条文番号を1行ずつ。確からしさを3段階で）

<<<条文・法令>>>
法令名 | 条 | どういう文脈で出たか
（この区間で挙がった法令名と条文番号を1行ずつ）

<<<公表状態>>>
主題 | 講師の発言をそのまま引く | 時刻
（★講師が「まだ公にしていない」と断った箇所を1行ずつ。無ければ「なし」）

■ ★「公表状態」について（★これは確信度とは別物）

★次のような箇所を拾う。講師が★確信して話していても、
  それが★まだ公にされていないなら、ここに入れる:
  ・「オフレコ」「ここだけの話」「フライング」
  ・「まだどこにも言っていない」「まだ公表していない」「未公表」
  ・「◯◯にもまだ伝えていない」「まだ発表されていない」
  ・「これは私見ですが」「政府見解ではありませんが」
  ・「後で立法担当解説に書かれる」「今後の国会審議で明らかにされる」
    （＝★いま時点では公的な裏付けが無い、という意味）
★上に挙げた言い方は例です。同じ趣旨の別の言い回しも拾ってください。
★拾ったものを本文に溶かし込まないこと。ここに分けて出す。
★これは「入れてはいけない」という意味ではありません。
  入れるかどうかは人が決めるので、★判断の材料として残してください。"""

_OVERVIEW_SYSTEM = (
    "あなたは記録係です。★章の見出しだけを見て、全体像を書きます。"
    "見出しに無いことを足しません。"
)

_ABOUT_SYSTEM = (
    "あなたは記録係です。★配付資料に書かれていることだけから、"
    "その講演の題名・日時・講師・主催を書き出します。"
    "★書かれていない項目は「不明」と書きます。推測しません。"
)

_NEXT_SYSTEM = (
    "あなたは行政書士事務所の記録係です。★章の見出しと条文の一覧だけを見て、"
    "一次資料に当たるべきものを挙げます。★講演で語られていない推測は書きません。"
)


def generate_chunk_chapters(chunk, i, n, material_text, api_key):
    """★1区画ぶんの章を作る。戻り値は RB.ChunkResult。"""
    client = OpenAI(api_key=api_key)
    mat = ""
    if material_text and material_text.strip():
        mat = ("\n【配付資料（用語・固有名詞の表記はこちらに合わせる）】\n"
               + material_text[:MAX_MATERIAL_CHARS] + "\n")
    corr = ""
    if CORRECTIONS:
        corr = ("   ★この講演でよく出る誤変換の対応表（★機械的に置換せず、"
                "文脈を見て判断すること）:\n"
                + "\n".join("     " + c for c in CORRECTIONS))
    prompt = _CHUNK_TEMPLATE.format(
        i=i + 1, n=n, context=chunk.context or "（なし）",
        body=chunk.body, material=mat, corrections=corr)
    try:
        with api_usage.record(app=USAGE_APP, site="report.chunk",
                              provider="openai", operation="chat.completions",
                              model="gpt-4o") as _rec:
            resp = client.chat.completions.create(
                model="gpt-4o",
                messages=[{"role": "system", "content": _CHUNK_SYSTEM},
                          {"role": "user", "content": prompt}],
                temperature=0.3,
                max_tokens=8000,
            )
            _rec.sdk_response(resp)
        ch = resp.choices[0]
        u = getattr(resp, "usage", None)
        return RB.ChunkResult(
            index=chunk.index, ok=True, text=ch.message.content or "",
            finish_reason=getattr(ch, "finish_reason", "") or "",
            in_tokens=getattr(u, "prompt_tokens", None),
            out_tokens=getattr(u, "completion_tokens", None))
    except Exception as e:                                   # noqa: BLE001
        # ★握り潰さない。失敗した区画は「捨てた量」に入り、レポートに出る。
        return RB.ChunkResult(index=chunk.index, ok=False,
                              error="%s: %s" % (type(e).__name__, str(e)[:200]))


def generate_overview(titles, api_key):
    """★材料は章の見出しだけ。本編を渡さない（もう一度要約されないため）。"""
    if not titles:
        return ""
    client = OpenAI(api_key=api_key)
    try:
        with api_usage.record(app=USAGE_APP, site="report.overview",
                              provider="openai", operation="chat.completions",
                              model="gpt-4o") as _rec:
            resp = client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": _OVERVIEW_SYSTEM},
                    {"role": "user", "content":
                        "次は、ある録音のレポートの章の見出し一覧です。"
                        "これだけを材料に、全体が何の話で、聞き手が何を"
                        "持ち帰るべきかを300字程度で書いてください。"
                        "★見出しに無いことは書かないでください。\n\n"
                        + "\n".join("- " + t for t in titles)},
                ],
                temperature=0.3,
                max_tokens=800,
            )
            _rec.sdk_response(resp)
        return resp.choices[0].message.content or ""
    except Exception as e:                                   # noqa: BLE001
        st.warning(f"⚠️ 全体像の生成に失敗しました（本編は影響を受けません）: {e}")
        return ""


def generate_about(material_text, file_labels, api_key):
    """★題名・日時・講師。材料は配付資料だけ（本編を渡さない）。"""
    if not (material_text or "").strip():
        return "題名／日時／講師：不明（配付資料が渡されていません）"
    client = OpenAI(api_key=api_key)
    try:
        with api_usage.record(app=USAGE_APP, site="report.about",
                              provider="openai", operation="chat.completions",
                              model="gpt-4o") as _rec:
            resp = client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": _ABOUT_SYSTEM},
                    {"role": "user", "content":
                        "次は、ある講演の配付資料から取り出した文字です。"
                        "題名／日時／講師／主催／配付資料 を箇条書きで"
                        "書き出してください。★書かれていない項目は「不明」と"
                        "書いてください。\n\n" + material_text[:MAX_MATERIAL_CHARS]},
                ],
                temperature=0.0,
                max_tokens=600,
            )
            _rec.sdk_response(resp)
        return resp.choices[0].message.content or "不明"
    except Exception as e:                                   # noqa: BLE001
        st.warning(f"⚠️ 「この記録について」の生成に失敗しました: {e}")
        return "不明（生成に失敗しました）"


def generate_next_steps(titles, laws, api_key):
    """★次に確かめること。材料は章の見出しと条文の一覧だけ。"""
    if not titles and not laws:
        return ""
    client = OpenAI(api_key=api_key)
    try:
        with api_usage.record(app=USAGE_APP, site="report.next_steps",
                              provider="openai", operation="chat.completions",
                              model="gpt-4o") as _rec:
            resp = client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": _NEXT_SYSTEM},
                    {"role": "user", "content":
                        "次は、ある講演のレポートの章の見出しと、"
                        "講演で挙がった条文・法令の一覧です。"
                        "事務所として一次資料（条文・通知・パブコメ等）に"
                        "当たるべきものを箇条書きで挙げてください。"
                        "★ここに無いことは書かないでください。\n\n"
                        "【章の見出し】\n" + "\n".join("- " + t for t in titles)
                        + "\n\n【条文・法令】\n"
                        + ("\n".join("- " + l for l in laws) or "（なし）")},
                ],
                temperature=0.2,
                max_tokens=1200,
            )
            _rec.sdk_response(resp)
        return resp.choices[0].message.content or ""
    except Exception as e:                                   # noqa: BLE001
        st.warning(f"⚠️ 「次に確かめること」の生成に失敗しました: {e}")
        return ""


def generate_report_full(transcript, file_labels, material_text, api_key,
                         segments=None, title="", allow_over_limit=False,
                         audio_notes=None):
    """★全文を区画に分けて読み、章を連結してレポートにする。

    ★ここで要約はしない。区画の出力をそのまま連結する。
    ★見込みが上限を超えたら止める。抜き取りには倒れない。
    """
    chunks = RB.split_for_reading(transcript)

    # ★資料を切るなら、切ったことを画面とレポートの両方に出す（黙って捨てない）。
    material_note = ""
    if material_text and len(material_text) > MAX_MATERIAL_CHARS:
        material_note = ("★配付資料は %s字あり、各区画には先頭 %s字だけを渡しています"
                         "（%s字は渡していません）。"
                         % (f"{len(material_text):,}", f"{MAX_MATERIAL_CHARS:,}",
                            f"{len(material_text) - MAX_MATERIAL_CHARS:,}"))
        st.warning("⚠️ " + material_note)

    est = RB.estimate_plan_jpy(
        total_chars=len(transcript),
        material_chars=len(material_text or ""),
        prompt_chars=len(_CHUNK_TEMPLATE),
        n_chunks=len(chunks),
        out_chars_per_chunk=int(RB.CHUNK_CHARS * 0.6))
    st.info("💴 " + RB.cost_note(est))
    if est["over_limit"] and not allow_over_limit:
        raise RB.CostLimitExceeded(
            "見込み %.0f円 が上限 %d円 を超えました（区画%d本）。"
            "★一部だけ読む形には倒しません。実行していません。"
            % (est["jpy"], est["limit_jpy"], len(chunks)))

    results = []
    pb = st.progress(0.0, text="区画 0/%d" % len(chunks))
    for i, c in enumerate(chunks):
        results.append(generate_chunk_chapters(c, i, len(chunks),
                                               material_text, api_key))
        pb.progress((i + 1) / len(chunks),
                    text="区画 %d/%d（%s字）" % (i + 1, len(chunks), f"{c.chars:,}"))

    cov = RB.coverage(len(transcript), chunks, results)
    bad = [r for r in results if not r.ok]
    if bad:
        st.error("★%d区画が読めませんでした。レポートに欠落として明記します。\n%s"
                 % (len(bad), "\n".join("- 区画%d: %s" % (r.index, r.error)
                                        for r in bad)))
    if cov["truncated"]:
        st.warning("★出力の上限で途中までになった区画: %s"
                   % ", ".join(str(i) for i in cov["truncated"]))

    # ★区画の出力を機械的に読み取る。一覧の統合に LLM を通さない。
    parsed = [RB.parse_chunk_output(r.text if r.ok else "") for r in results]
    chapters = [p["chapters"] for p, r in zip(parsed, results)
                if r.ok and p["chapters"].strip()]
    no_heading = RB.chapters_missing_heading(parsed, results)
    if no_heading:
        # ★破綻点の検知。連結すると見た目は通るので、ここで必ず出す。
        st.error("★章の形になっていない区画があります: %s"
                 % ", ".join("区画%d" % i for i in no_heading))

    terms = RB.merge_rows(parsed, "terms")
    # ★「直す前」が元の文字起こしに無い訂正を落とす。対応表を書き写しただけの
    #   行が混ざると、訂正の件数が嘘になる（2026-08-21 実測で5組あった）。
    terms, dropped_terms = RB.filter_corrections_by_source(terms, transcript)
    if dropped_terms:
        st.warning("★実際には直していない訂正を %d件 落としました（元の文字起こしに"
                   "その語が出てきません）: %s"
                   % (len(dropped_terms), " / ".join(dropped_terms[:5])))
    numbers = RB.merge_rows(parsed, "numbers")
    laws = RB.merge_rows(parsed, "laws")
    # ★公表状態は本文に混ぜず、判断の対象として別の節に出す。
    publicity = RB.merge_rows(parsed, "publicity")
    if publicity:
        st.warning("★公的な裏付けが無い箇所が %d件 あります。"
                   "レポートの「★公的な裏付けが無い箇所」を見て、"
                   "知識ベースに入れるかを決めてください。" % len(publicity))
    titles = RB.chapter_titles(chapters)

    # ★本文に数値が残っていないかを機械的に数える。指示だけでは守られない
    #   ことが 2026-08-21 の実測で分かっているので、出た物を表に出す。
    num_found = RB.numbers_in_body("\n".join(chapters))
    if num_found["total"]:
        st.warning("★本文に数値が %d行 残っています（日付%d・条文番号%d・時刻%d）。"
                   "レポートの「本文の数値の点検」に出しています。"
                   % (num_found["total"], len(num_found["date"]),
                      len(num_found["article"]), len(num_found["timestamp"])))

    about = generate_about(material_text, file_labels, api_key)
    overview = generate_overview(titles, api_key)
    next_steps = generate_next_steps(titles, laws, api_key)
    gaps = RB.silent_gaps(segments) if segments else []
    note = "\n\n".join([x for x in (audio_notes or []) if x]
                        + ([material_note] if material_note else []))
    return RB.assemble_full(about, overview, chapters, terms, numbers, laws,
                            cov, gaps, next_steps, note, no_heading,
                            num_found, publicity), cov, est


# ═══════════════════════════════════════════
# GPT：構造化サマリー（JSON）
# ═══════════════════════════════════════════
def generate_summary_json(chapter_titles, report, material_text, api_key):
    """★2026-08-20 変更。文字起こしの先頭6000字を渡すのをやめた。

    先頭だけを見て「全体のサマリー」を名乗るのは、★中身が落ちたことが
    分からない形（今回直したのと同じ形）。ここでは
    ・章の見出し（全章ぶん・落ちない）
    ・本編そのもの
    を渡す。★本編が長くて入らない場合は、切ったことを本文に書く。
    """
    client = OpenAI(api_key=api_key)

    titles_txt = "\n".join("- " + t for t in (chapter_titles or [])) or "（なし）"
    safe_report = report
    cut_note = ""
    if len(report) > 60000:
        safe_report = report[:60000]
        cut_note = ("\n★注意: レポートが長いため先頭60,000字だけを渡しています"
                    "（元 %s字）。このサマリーは全体を見ていません。\n"
                    % f"{len(report):,}")
    mat_note = "補足資料の情報も反映してください。" if material_text else ""

    prompt = f"""以下のレポートから、構造化サマリーをJSON形式で作成してください。{mat_note}
{cut_note}
【章の見出し（★全章ぶん）】
{titles_txt}

【レポート本文】
{safe_report}

以下のJSON構造で出力してください（日本語で）。
コードブロック（```）は使わず、JSONのみを出力してください。

{{
  "title": "会議・メモのタイトル（15〜30文字）",
  "date": "推定日付または「不明」",
  "type": "会議 / 1on1 / ブレスト / 講義 / その他",
  "duration": "推定XX分",
  "urgency": "高 / 中 / 低",
  "one_line": "この会議・メモを一文で表すと（30〜50文字）",
  "participants": ["参加者1", "参加者2"],
  "flow": [
    {{"time": "序盤", "topic": "トピック名", "summary": "内容の要約（30〜60文字）"}},
    {{"time": "中盤", "topic": "トピック名", "summary": "内容の要約（30〜60文字）"}},
    {{"time": "終盤", "topic": "トピック名", "summary": "内容の要約（30〜60文字）"}}
  ],
  "decisions": [
    {{"title": "決定事項名", "detail": "詳細説明"}}
  ],
  "actions": [
    {{"priority": "高/中/低", "who": "担当者", "what": "タスク内容", "when": "期限"}}
  ],
  "concerns": [
    {{"title": "懸念・リスク名", "detail": "詳細説明"}}
  ],
  "next_topics": ["次回以降の検討事項1", "次回以降の検討事項2"],
  "key_numbers": [
    {{"label": "指標名・数値名", "value": "具体的な数値・データ"}}
  ],
  "keywords": ["重要キーワード1", "重要キーワード2", "重要キーワード3", "重要キーワード4", "重要キーワード5"]
}}

注意：
- decisionsは実際に決定したことのみ。なければ空配列[]
- actionsは具体的なタスク。なければ空配列[]
- concernsはリスク・懸念・未解決事項。なければ空配列[]
- key_numbersは具体的な数値が言及された場合のみ。なければ空配列[]
"""

    try:
        resp = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "あなたは会議の内容を正確に構造化するアナリストです。指示通りのJSONのみを出力します。"},
                {"role": "user", "content": prompt}
            ],
            temperature=0.2,
            max_tokens=2000,
            response_format={"type": "json_object"}
        )
        return json.loads(resp.choices[0].message.content)
    except Exception as e:
        st.error(f"構造化サマリー生成エラー: {e}")
        return None


# ═══════════════════════════════════════════
# 構造化サマリー → HTML
# ═══════════════════════════════════════════
def summary_to_html(data, file_labels, generated_at):
    urgency_color = {"高": "#e53e5a", "中": "#f5a623", "低": "#22c38e"}.get(data.get("urgency", "中"), "#888")
    urgency_bg    = {"高": "#fff0f2", "中": "#fff8ee", "低": "#f0fff8"}.get(data.get("urgency", "中"), "#f5f5f5")

    flow_items = data.get("flow", [])
    flow_html = ""
    for i, f in enumerate(flow_items):
        connector = '<div class="flow-arrow">↓</div>' if i < len(flow_items) - 1 else ""
        flow_html += f"""
        <div class="flow-item">
          <div class="flow-time">{f.get('time','')}</div>
          <div class="flow-content">
            <div class="flow-topic">{f.get('topic','')}</div>
            <div class="flow-summary">{f.get('summary','')}</div>
          </div>
        </div>{connector}"""

    decisions = data.get("decisions", [])
    dec_html = "".join(
        f'<div class="card-item card-decision"><div class="card-item-title">✅ {d.get("title","")}</div><div class="card-item-detail">{d.get("detail","")}</div></div>'
        for d in decisions
    ) if decisions else '<div class="empty-note">言及なし</div>'

    actions = data.get("actions", [])
    pc_map = {"高": "#e53e5a", "中": "#f5a623", "低": "#22c38e"}
    act_html = "".join(
        f'''<div class="action-row">
          <span class="action-priority" style="background:{pc_map.get(a.get("priority","中"),"#888")}20;color:{pc_map.get(a.get("priority","中"),"#888")};border:1px solid {pc_map.get(a.get("priority","中"),"#888")}40">{a.get("priority","")}</span>
          <div class="action-body">
            <div class="action-what">{a.get("what","")}</div>
            <div class="action-meta">👤 {a.get("who","未定")} &nbsp;｜&nbsp; 📅 {a.get("when","期限未定")}</div>
          </div>
        </div>'''
        for a in sorted(actions, key=lambda x: {"高":0,"中":1,"低":2}.get(x.get("priority","中"),1))
    ) if actions else '<div class="empty-note">言及なし</div>'

    concerns = data.get("concerns", [])
    con_html = "".join(
        f'<div class="card-item card-concern"><div class="card-item-title">⚠️ {c.get("title","")}</div><div class="card-item-detail">{c.get("detail","")}</div></div>'
        for c in concerns
    ) if concerns else '<div class="empty-note">言及なし</div>'

    nexts = data.get("next_topics", [])
    next_html = "".join(f"<li>{n}</li>" for n in nexts) if nexts else '<li class="empty-note">言及なし</li>'

    nums = data.get("key_numbers", [])
    num_html = "".join(
        f'<div class="kpi-card"><div class="kpi-value">{n.get("value","")}</div><div class="kpi-label">{n.get("label","")}</div></div>'
        for n in nums
    )

    keywords = data.get("keywords", [])
    kw_html = "".join(f'<span class="keyword">{k}</span>' for k in keywords)
    participants = data.get("participants", [])
    par_html = "・".join(participants) if participants else "不明"
    files_html = "・".join(file_labels)

    return f"""<!DOCTYPE html>
<html lang="ja">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>{data.get('title','構造化サマリー')}</title>
<style>
@import url('https://fonts.googleapis.com/css2?family=Noto+Sans+JP:wght@300;400;500;700;900&display=swap');
:root {{
  --ink:#1a2140;--ink2:#4a567a;--ink3:#8892b0;
  --line:#e2e8f0;--bg:#f8faff;--card:#ffffff;
  --blue:#3b6ef0;--blue-lt:#eef2ff;
  --green:#22c38e;--red:#e53e5a;--amber:#f5a623;
  --radius:10px;--shadow:0 2px 12px rgba(26,33,64,.08);
}}
*{{box-sizing:border-box;margin:0;padding:0;}}
body{{font-family:'Noto Sans JP',sans-serif;background:var(--bg);color:var(--ink);font-size:14px;line-height:1.7;}}
.page{{max-width:860px;margin:0 auto;padding:40px 32px 80px;}}
.doc-header{{border-bottom:3px solid var(--blue);padding-bottom:24px;margin-bottom:32px;}}
.doc-meta{{display:flex;gap:10px;flex-wrap:wrap;margin-bottom:12px;}}
.meta-chip{{display:inline-flex;align-items:center;gap:5px;font-size:11px;font-weight:600;letter-spacing:.06em;color:var(--ink3);background:white;border:1px solid var(--line);border-radius:99px;padding:3px 11px;}}
.doc-title{{font-size:clamp(20px,3vw,28px);font-weight:900;color:var(--ink);line-height:1.3;margin-bottom:12px;}}
.one-line{{font-size:14px;color:var(--ink2);background:var(--blue-lt);border-left:4px solid var(--blue);padding:10px 16px;border-radius:0 8px 8px 0;font-weight:500;}}
.urgency-badge{{display:inline-flex;align-items:center;gap:5px;font-size:12px;font-weight:700;padding:4px 14px;border-radius:99px;background:{urgency_bg};color:{urgency_color};border:1.5px solid {urgency_color}40;margin-top:12px;}}
.kpi-row{{display:flex;gap:12px;flex-wrap:wrap;margin-bottom:32px;}}
.kpi-card{{background:var(--card);border:1px solid var(--line);border-radius:var(--radius);padding:16px 20px;min-width:120px;text-align:center;box-shadow:var(--shadow);flex:1;}}
.kpi-value{{font-size:22px;font-weight:900;color:var(--blue);line-height:1.2;}}
.kpi-label{{font-size:11px;color:var(--ink3);margin-top:4px;}}
.section{{margin-bottom:32px;}}
.section-title{{font-size:11px;font-weight:800;letter-spacing:.14em;text-transform:uppercase;color:var(--blue);margin-bottom:12px;display:flex;align-items:center;gap:8px;}}
.section-title::after{{content:'';flex:1;height:1px;background:var(--line);}}
.flow-wrap{{display:flex;flex-direction:column;}}
.flow-item{{display:flex;gap:16px;align-items:flex-start;background:var(--card);border:1px solid var(--line);border-radius:var(--radius);padding:14px 18px;width:100%;box-shadow:var(--shadow);}}
.flow-arrow{{text-align:center;color:var(--ink3);font-size:18px;padding:4px 0;}}
.flow-time{{font-size:11px;font-weight:700;color:var(--blue);background:var(--blue-lt);padding:3px 10px;border-radius:99px;white-space:nowrap;flex-shrink:0;align-self:flex-start;margin-top:2px;}}
.flow-topic{{font-size:14px;font-weight:700;color:var(--ink);margin-bottom:4px;}}
.flow-summary{{font-size:13px;color:var(--ink2);}}
.card-item{{background:var(--card);border-radius:var(--radius);padding:14px 18px;margin-bottom:10px;border:1px solid var(--line);box-shadow:var(--shadow);}}
.card-decision{{border-left:4px solid var(--green);}}
.card-concern{{border-left:4px solid var(--amber);}}
.card-item-title{{font-size:14px;font-weight:700;color:var(--ink);margin-bottom:4px;}}
.card-item-detail{{font-size:13px;color:var(--ink2);}}
.action-row{{display:flex;gap:14px;align-items:flex-start;background:var(--card);border:1px solid var(--line);border-radius:var(--radius);padding:12px 16px;margin-bottom:8px;box-shadow:var(--shadow);}}
.action-priority{{font-size:11px;font-weight:700;padding:3px 10px;border-radius:99px;white-space:nowrap;flex-shrink:0;}}
.action-what{{font-size:14px;font-weight:600;color:var(--ink);margin-bottom:4px;}}
.action-meta{{font-size:12px;color:var(--ink3);}}
.next-list{{list-style:none;display:flex;flex-direction:column;gap:8px;}}
.next-list li{{background:var(--card);border:1px solid var(--line);border-radius:var(--radius);padding:10px 16px;font-size:13px;color:var(--ink2);box-shadow:var(--shadow);}}
.next-list li::before{{content:"→ ";color:var(--blue);font-weight:700;}}
.keyword-wrap{{display:flex;flex-wrap:wrap;gap:8px;}}
.keyword{{background:var(--blue-lt);color:var(--blue);border:1px solid #c0cef8;border-radius:99px;padding:4px 14px;font-size:12px;font-weight:600;}}
.two-col{{display:grid;grid-template-columns:1fr 1fr;gap:24px;}}
.empty-note{{font-size:13px;color:var(--ink3);font-style:italic;padding:8px 4px;}}
.doc-footer{{margin-top:48px;padding-top:16px;border-top:1px solid var(--line);font-size:11px;color:var(--ink3);}}
.files-note{{font-size:12px;color:var(--ink3);margin-top:4px;}}
@media(max-width:600px){{.page{{padding:24px 16px 60px;}}.two-col{{grid-template-columns:1fr;}}}}
@media print{{body{{background:white;}}.page{{padding:20px;max-width:100%;}}.card-item,.action-row,.flow-item,.next-list li{{-webkit-print-color-adjust:exact;print-color-adjust:exact;break-inside:avoid;}}.section{{break-inside:avoid;}}}}
</style>
</head>
<body>
<div class="page">
  <div class="doc-header">
    <div class="doc-meta">
      <span class="meta-chip">📅 {data.get('date','不明')}</span>
      <span class="meta-chip">🎙️ {data.get('type','会議')}</span>
      <span class="meta-chip">⏱ {data.get('duration','不明')}</span>
      <span class="meta-chip">👥 {par_html}</span>
    </div>
    <div class="doc-title">{data.get('title','構造化サマリー')}</div>
    <div class="one-line">{data.get('one_line','')}</div>
    <div class="urgency-badge">{'🔴' if data.get('urgency')=='高' else '🟡' if data.get('urgency')=='中' else '🟢'} 緊急度：{data.get('urgency','中')}</div>
    <div class="files-note">📁 対象ファイル：{files_html}</div>
  </div>
  {"<div class='kpi-row'>" + num_html + "</div>" if nums else ""}
  <div class="section">
    <div class="section-title">📋 話の流れ・構成</div>
    <div class="flow-wrap">{flow_html}</div>
  </div>
  <div class="two-col">
    <div class="section">
      <div class="section-title">✅ 決定事項</div>{dec_html}
    </div>
    <div class="section">
      <div class="section-title">⚠️ 懸念・リスク</div>{con_html}
    </div>
  </div>
  <div class="section">
    <div class="section-title">🎯 アクションアイテム</div>{act_html}
  </div>
  <div class="section">
    <div class="section-title">🔄 次回以降の検討事項</div>
    <ul class="next-list">{next_html}</ul>
  </div>
  {"<div class='section'><div class='section-title'>🏷 キーワード</div><div class='keyword-wrap'>" + kw_html + "</div></div>" if keywords else ""}
  <div class="doc-footer">
    <div>📁 {files_html}</div>
    <div>🕐 生成日時：{generated_at}</div>
  </div>
</div>
</body>
</html>"""


# ═══════════════════════════════════════════
# UI：サイドバー
# ═══════════════════════════════════════════
with st.sidebar:
    st.header("⚙️ 設定")
    if st.session_state.api_key:
        st.success("✓ OpenAI APIキー設定済み")
    else:
        api_key_input = st.text_input(
            "OpenAI APIキー", type="password",
            placeholder="sk-...",
            help="環境変数 OPENAI_API_KEY が未設定の場合に入力"
        )
        if api_key_input:
            st.session_state.api_key = api_key_input
            st.success("✓ APIキー設定済み")

    if NOTION_API_KEY:
        st.success("✓ Notion APIキー設定済み")
    else:
        st.warning("⚠️ NOTION_TOKEN 未設定（Notion保存不可）")

    st.divider()
    st.markdown(r"""
### 📋 入力ソース
| | 対応 |
|---|---|
| 🎵 音声 | MP3/WAV/M4A/WebM |
| 🎬 YouTube | URLから字幕取得 |
| 📝 テキスト | TXT/MD直接入力 |

**📄 補足資料（任意）**
PDF / PPTX / DOCX

### 📤 出力
- PLAUDレポート（.md）
- マインドマップ（Markmap）
- Notion研修DB保存

### ⚙️ 動かす場所
★PC-B のローカル（2026-08-20 に外部クラウドから移設）。
起動は他のローカルアプリと同じ `C:\dev\start_business.bat`
（このアプリは http://localhost:8512）。

### 🔑 キーの置き場
★`C:\Users\ohga\OneDrive\secrets\.env` の1か所だけ。
shared-lib/env_loader が読むので★画面での入力は不要です。
（下の入力欄は .env が読めないときの最後の手段）
""")

    # ── ★語彙ヒントの状態（黙って捨てない）──
    if VOCAB_TERMS:
        st.caption("🗣️ 語彙ヒント: %d語中 %d語を音声認識に渡します"
                   % (len(VOCAB_TERMS), VOCAB_USED))
        if VOCAB_DROPPED:
            st.warning("★%d語が上限（%d字）に入りませんでした。"
                       "vocab/legal_ja.txt の上の方に並べ替えてください。"
                       % (VOCAB_DROPPED, RB.VOCAB_PROMPT_MAX_CHARS))
    else:
        st.caption("🗣️ 語彙ヒントなし（vocab/legal_ja.txt が読めません）")

    # ── ★費用の上限を外す（日常の操作と分ける）──
    #   ★既定は外れている。押しやすい場所には置かない。
    #   上限を超えたときに「一部だけ読む」へ倒れないための、唯一の逃げ道。
    with st.expander("★費用の上限を外す（%d円）" % RB.COST_LIMIT_JPY):
        st.caption(
            "1回あたりの見込みが %d円 を超えると、実行せずに止まります。"
            "★止まったときに一部だけ読む形には倒れません。"
            "承知のうえで続ける場合だけ、ここに印を付けてください。"
            % RB.COST_LIMIT_JPY)
        st.session_state["allow_over_limit"] = st.checkbox(
            "上限を超えても実行する", value=False, key="allow_over_limit_cb")
        if st.session_state.get("allow_over_limit"):
            st.warning("★上限が外れています。実行前に見込み額を確かめてください。")


# ═══════════════════════════════════════════
# UI：メイン
# ═══════════════════════════════════════════
st.title("🎙️ 音声メモアプリ Pro")
st.caption("音声 / YouTube / テキスト → PLAUDレポート ／ マインドマップ ／ Notion保存")

if not st.session_state.api_key:
    st.warning("⚠️ OpenAI APIキーが未設定です。環境変数 OPENAI_API_KEY を設定するか、サイドバーで入力してください。")
    st.stop()

# ── 入力ソース選択 ──
st.subheader("① 入力ソースを選択")
source_type = st.radio(
    "入力ソース",
    ["🎵 音声ファイル", "🎬 YouTube URL", "📝 テキストファイル"],
    horizontal=True,
    label_visibility="collapsed",
)

audio_files = []
youtube_url = ""
text_files = []
pasted_text = ""

if source_type == "🎵 音声ファイル":
    st.caption("複数ファイルはファイル名順で結合し、**1つのレポート**を作成します。")
    audio_files = st.file_uploader(
        "MP3・WAV・M4A・WebM",
        type=["mp3", "wav", "m4a", "webm"],
        accept_multiple_files=True,
    )

elif source_type == "🎬 YouTube URL":
    st.caption("公開動画で字幕（自動生成を含む）が有効なものに対応します。")
    youtube_url = st.text_input(
        "YouTube URL", placeholder="https://www.youtube.com/watch?v=..."
    )

elif source_type == "📝 テキストファイル":
    st.caption("TXT / MD ファイルをアップロード、またはテキストを直接貼り付けてください。Whisperをスキップして直接レポート生成します。")
    text_files = st.file_uploader(
        "TXT・MD（ファイルアップロード）",
        type=["txt", "md"],
        accept_multiple_files=True,
    )
    st.caption("または")
    pasted_text = st.text_area(
        "テキスト直接貼り付け",
        height=300,
        placeholder="文字起こしテキストをここに貼り付けてください...",
        key="pasted_text_input",
        label_visibility="collapsed",
    )

st.markdown("---")

# ── 補足資料 ──
st.subheader("② 補足資料（任意・複数可）")
st.caption("会議資料・スライドなど。なくても動作します。")
material_files = st.file_uploader(
    "PDF・PPTX・DOCX",
    type=["pdf", "pptx", "ppt", "docx", "doc"],
    accept_multiple_files=True,
    key="material_uploader",
)

st.markdown("---")

# ── 添付資料（Notion保存用）──
st.subheader("③ 添付資料（Notion保存用・任意・複数可）")
st.caption("Notionページに添付ファイル一覧として記録します（ファイル名と容量のみ保存）。")
attachment_files = st.file_uploader(
    "PDF・PPTX・DOCX",
    type=["pdf", "pptx", "ppt", "docx", "doc"],
    accept_multiple_files=True,
    key="attachment_uploader",
)

# ── 入力確認表示 ──
has_input = bool(audio_files or youtube_url.strip() or text_files or pasted_text.strip())

if has_input:
    if audio_files:
        def sort_key(f):
            nums = re.findall(r'\d+', f.name)
            return "".join(nums).zfill(20) if nums else f.name
        sorted_audio = sorted(audio_files, key=sort_key)

        with st.expander(f"📋 音声ファイル {len(sorted_audio)}件", expanded=True):
            for i, f in enumerate(sorted_audio, 1):
                mb = f.size / (1024 * 1024)
                c1, c2, c3 = st.columns([5, 2, 2])
                c1.write(f"**{i}.** {f.name}")
                c2.caption(f"{mb:.1f} MB")
                c3.caption("🔧 要圧縮" if mb > 24 else "✅ OK")
            if len(sorted_audio) > 1:
                st.info(f"💡 {len(sorted_audio)}件を順番に文字起こし → 結合 → 1本のレポートを生成します。")
    elif youtube_url.strip():
        st.info(f"🎬 YouTube: {youtube_url.strip()}")
    elif pasted_text.strip():
        preview = pasted_text.strip()[:80].replace("\n", " ")
        st.info(f"📝 貼り付けテキスト（{len(pasted_text.strip()):,}文字）: {preview}...")
    elif text_files:
        st.info(f"📝 テキストファイル: {', '.join(f.name for f in text_files)}")

    if material_files:
        st.info(f"📎 補足資料（{len(material_files)}件）: {', '.join(f.name for f in material_files)}")

    st.markdown("---")
    btn_col1, btn_col2 = st.columns(2)
    with btn_col1:
        if st.button("🗑️ 処理結果をクリア", use_container_width=True):
            st.session_state.results = []
            st.rerun()
    with btn_col2:
        run = st.button("🚀 処理開始", type="primary", use_container_width=True)

    if run:
        # ── 補足資料抽出 ──
        combined_material = None
        if material_files:
            with st.spinner("📄 補足資料を読み込み中..."):
                mat_texts = []
                for mf in material_files:
                    t = extract_material_text(mf)
                    if t and not t.startswith("["):
                        mat_texts.append(f"=== {mf.name} ===\n{t}")
                    else:
                        st.warning(f"⚠️ {mf.name}: {t}")
            if mat_texts:
                combined_material = "\n\n".join(mat_texts)
                st.success(f"✅ 資料 {len(mat_texts)}件 読み込み完了")

        # ── STEP 1：テキスト取得 ──
        st.markdown("---")
        raw_transcript = ""
        file_labels = []
        transcripts_per_file = {}
        # ★時刻付きセグメント。音声以外の入力では空のまま（無い物を作らない）。
        all_segments = []
        # ★音質を落としたなど、レポートに残すべき経過。
        audio_notes = []
        source_label = "音声"  # Notion保存用

        if source_type == "🎵 音声ファイル":
            source_label = "音声"
            st.markdown("### 🎧 STEP1：文字起こし")
            all_tmp_paths = []
            for idx, audio_file in enumerate(sorted_audio):
                st.markdown(f"**[{idx+1}/{len(sorted_audio)}]** {audio_file.name}")
                suffix = Path(audio_file.name).suffix
                with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as f:
                    f.write(audio_file.read())
                    tmp_path = f.name
                    all_tmp_paths.append(tmp_path)
                with st.spinner("  文字起こし中..."):
                    # ★文字起こしの分も「1回あたり」に数える（2026-08-20）。
                    #   ここを数えないと、関所が費用の一部しか見ないことになる。
                    mins = _audio_minutes(tmp_path)
                    if mins:
                        we = RB.estimate_whisper_jpy(mins)
                        st.info("🎧 文字起こしの見込み: %.0f円（%.1f分・$%.2f）"
                                % (we["jpy"], we["minutes"], we["usd"]))
                        if (we["over_limit"]
                                and not st.session_state.get("allow_over_limit")):
                            st.error("🛑 文字起こしだけで見込み %.0f円 が上限 %d円 を"
                                     "超えました。★実行していません。"
                                     % (we["jpy"], we["limit_jpy"]))
                            st.stop()
                    tr, segs, tinfo = transcribe_audio(
                        tmp_path, st.session_state.api_key, VOCAB_PROMPT)
                if tr:
                    all_segments.extend(segs or [])
                    if tinfo.get("note"):
                        audio_notes.append(tinfo["note"])
                    transcripts_per_file[audio_file.name] = tr
                    st.success(f"  ✅ 完了（{len(tr):,}文字）")
                else:
                    st.error("  ❌ 失敗")
            for p in all_tmp_paths:
                if os.path.exists(p):
                    os.remove(p)
            if not transcripts_per_file:
                st.error("文字起こしに成功したファイルがありません。")
                st.stop()
            file_labels = list(transcripts_per_file.keys())
            raw_transcript = (
                list(transcripts_per_file.values())[0]
                if len(file_labels) == 1
                else "\n\n".join(f"--- {k} ---\n{v}" for k, v in transcripts_per_file.items())
            )
            st.success(f"✅ 文字起こし完了（合計 {len(raw_transcript):,}文字）")

        elif source_type == "🎬 YouTube URL":
            source_label = "YouTube"
            st.markdown("### 🎬 STEP1：YouTube字幕取得")
            with st.spinner("字幕を取得中..."):
                raw_transcript, video_id = get_youtube_transcript(youtube_url.strip())
            if not raw_transcript:
                st.error(f"❌ 取得失敗")
                # video_id には失敗理由が入っている
                for line in video_id.splitlines():
                    st.error(line)
                st.info(
                    "💡 **対処法**\n\n"
                    "**① 字幕をコピーして貼り付け（最も簡単）**\n"
                    "  1. YouTube動画を開き、動画下の「…」→「文字起こし」をクリック\n"
                    "  2. 表示されたテキストを全選択してコピー\n"
                    "  3. 入力ソースを **「📝 テキストファイル」** に切り替えてテキスト欄に貼り付け\n\n"
                    "**② yt-dlp フォールバック（音声からWhisper文字起こし）**\n"
                    "  - `pip install yt-dlp` でインストールすると字幕なし動画にも対応します\n\n"
                    "**③ その他の原因**\n"
                    "  - IP ブロック → しばらく待つか別のネットワークで試してください\n"
                    "  - 年齢制限 → 現時点では取得不可（API 制限）\n"
                    "  - URL フォーマット → `watch?v=` / `youtu.be/` / Shorts URL に対応しています"
                )
                st.stop()
            file_labels = [youtube_url.strip()]
            transcripts_per_file = {youtube_url.strip(): raw_transcript}
            st.success(f"✅ 字幕取得完了（{len(raw_transcript):,}文字）")

        elif source_type == "📝 テキストファイル":
            source_label = "テキスト"
            st.markdown("### 📝 STEP1：テキスト読み込み")
            # 貼り付けテキストを優先。なければファイルを読み込む
            if pasted_text.strip():
                content = pasted_text.strip()
                transcripts_per_file["貼り付けテキスト"] = content
                file_labels.append("貼り付けテキスト")
                st.success(f"✅ 貼り付けテキスト（{len(content):,}文字）")
                if text_files:
                    st.info("💡 ファイルと貼り付けテキストの両方が入力されています。貼り付けテキストを優先して処理します。")
            else:
                for tf in text_files:
                    content = tf.read().decode("utf-8", errors="replace")
                    transcripts_per_file[tf.name] = content
                    file_labels.append(tf.name)
                    st.success(f"✅ {tf.name}（{len(content):,}文字）")
            if not transcripts_per_file:
                st.error("テキストが入力されていません。")
                st.stop()
            raw_transcript = "\n\n".join(
                f"--- {k} ---\n{v}" for k, v in transcripts_per_file.items()
            ) if len(transcripts_per_file) > 1 else list(transcripts_per_file.values())[0]

        # ── ★全文を読む（抜き取りはしない）──
        combined_transcript = raw_transcript

        st.markdown("### 📊 STEP2：レポート生成（★全文を読みます）")
        try:
            with st.spinner("区画ごとに読んでいます..."):
                report, coverage_info, cost_info = generate_report_full(
                    combined_transcript, file_labels, combined_material,
                    st.session_state.api_key, segments=all_segments,
                    audio_notes=audio_notes,
                    allow_over_limit=st.session_state.get("allow_over_limit", False),
                )
        except RB.CostLimitExceeded as e:
            # ★ここで一部だけ読む形に倒れない。止まる。
            st.error("🛑 %s" % e)
            st.info("続ける場合は、サイドバー最下部の「★費用の上限を外す」を"
                    "開いて印を付けてから、もう一度実行してください。")
            st.stop()
        if not report:
            st.error("レポート生成に失敗しました。")
            st.stop()
        st.success("✅ レポート完了（★%s字中 %s字を読了・%.1f%%）%s"
                   % (f"{coverage_info['total_chars']:,}",
                      f"{coverage_info['read_chars']:,}",
                      coverage_info["ratio"] * 100,
                      "（資料補完あり）" if combined_material else ""))

        # ── STEP 3：Markmap生成 ──
        st.markdown("### 🗺️ STEP3：マインドマップ生成")
        with st.spinner("Markmap生成中..."):
            markmap_md = generate_markmap(report, st.session_state.api_key)
        if markmap_md:
            st.success("✅ マインドマップ生成完了")
        else:
            st.warning("⚠️ マインドマップ生成に失敗しました")

        # ── STEP 4：構造化サマリー生成 ──
        st.markdown("### 📋 STEP4：構造化サマリー生成")
        with st.spinner("構造化サマリー生成中..."):
            summary_data = generate_summary_json(
                RB.chapter_titles([report]), report, combined_material,
                st.session_state.api_key
            )
        summary_html = None
        if summary_data:
            generated_at = datetime.now().strftime("%Y-%m-%d %H:%M")
            summary_html = summary_to_html(summary_data, file_labels, generated_at)
            st.success("✅ 構造化サマリー完了")

        # 結果保存
        result = {
            "label": "・".join(file_labels),
            "file_labels": file_labels,
            "date": datetime.now().strftime("%Y-%m-%d %H:%M"),
            "transcripts_per_file": transcripts_per_file,
            "combined_transcript": raw_transcript,
            "report": report,
            "markmap_md": markmap_md,
            "summary_html": summary_html,
            "summary_data": summary_data,
            "coverage": coverage_info,
            "cost_estimate": cost_info,
            "segments_count": len(all_segments),
            "has_material": combined_material is not None,
            "source_label": source_label,
            "youtube_url": youtube_url.strip() if source_type == "🎬 YouTube URL" else "",
            "attachment_file_info": [
                {"name": f.name, "size": f.size} for f in (attachment_files or [])
            ],
        }
        st.session_state.results = [result] + st.session_state.results
        st.balloons()
        st.success("🎉 処理完了！")


# ── 結果表示 ──
if st.session_state.results:
    st.markdown("---")
    hcol1, hcol2 = st.columns([4, 1])
    hcol1.header(f"📋 処理結果（{len(st.session_state.results)}件）")
    if hcol2.button("🗑️ 全クリア", key="clear_top"):
        st.session_state.results = []
        st.rerun()

    for result in st.session_state.results:
        mat_badge = "  📎 資料補完あり" if result["has_material"] else ""
        n_files = len(result["file_labels"])
        header_label = (
            f"📁 [{n_files}件統合] {result['label']}  —  {result['date']}{mat_badge}"
            if n_files > 1
            else f"📁 {result['label']}  —  {result['date']}{mat_badge}"
        )

        with st.expander(header_label, expanded=True):
            tab_labels = ["📄 文字起こし", "📊 レポート"]
            if result.get("markmap_md"):
                tab_labels.append("🗺️ マインドマップ")
            if result.get("summary_html"):
                tab_labels.append("📋 構造化サマリー")
            tabs = st.tabs(tab_labels)

            # 文字起こし
            with tabs[0]:
                if n_files > 1:
                    sub_labels = list(result["transcripts_per_file"].keys()) + ["📄 全文（結合）"]
                    sub_tabs = st.tabs(sub_labels)
                    for i, (fname, tr) in enumerate(result["transcripts_per_file"].items()):
                        with sub_tabs[i]:
                            st.text_area("", tr, height=220,
                                         key=f"tr_{result['date']}_{fname}")
                            st.download_button(
                                f"📥 {fname} (.txt)",
                                tr,
                                file_name=f"transcript_{Path(fname).stem}.txt",
                                mime="text/plain",
                                key=f"dtr_{result['date']}_{fname}"
                            )
                    with sub_tabs[-1]:
                        st.text_area("", result["combined_transcript"], height=300,
                                     key=f"tr_all_{result['date']}")
                        st.download_button(
                            "📥 全文字起こし（結合）(.txt)",
                            result["combined_transcript"],
                            file_name=f"transcript_combined_{result['date'].replace(':','').replace(' ','_')}.txt",
                            mime="text/plain",
                            key=f"dtr_all_{result['date']}"
                        )
                else:
                    fname = result["file_labels"][0]
                    tr = result["transcripts_per_file"][fname]
                    st.text_area("", tr, height=250,
                                 key=f"tr_{result['date']}_{fname}")
                    st.download_button(
                        "📥 文字起こし (.txt)",
                        tr,
                        file_name=f"transcript_{Path(fname).stem}.txt",
                        mime="text/plain",
                        key=f"dtr_{result['date']}_{fname}"
                    )

            # レポート
            with tabs[1]:
                if result["report"]:
                    st.markdown(result["report"])
                    fname_base = (
                        Path(result["file_labels"][0]).stem if n_files == 1
                        else f"combined_{result['date'].replace(':','').replace(' ','_')}"
                    )
                    dl_col, notion_col = st.columns([1, 1])
                    with dl_col:
                        st.download_button(
                            "📥 レポート (.md)",
                            result["report"],
                            file_name=f"report_{fname_base}.md",
                            mime="text/markdown",
                            key=f"drp_{result['date']}",
                        )
                    with notion_col:
                        if st.button("☁️ Notionに保存", key=f"notion_{result['date']}",
                                     disabled=not NOTION_API_KEY):
                            title = extract_title_from_report(result["report"])
                            tags  = extract_tags_from_report(result["report"])
                            summary_text = "\n".join(result["report"].splitlines()[:10])
                            # source_info: ファイル名またはURL一覧
                            src_info = result.get("file_labels", [])
                            if result.get("youtube_url"):
                                src_info = [result["youtube_url"]]
                            save_to_notion_kenshu(
                                title=title,
                                tags=tags,
                                source_type=result.get("source_label", "音声"),
                                report=result["report"],
                                summary=summary_text,
                                transcript=result.get("combined_transcript", ""),
                                markmap_md=result.get("markmap_md", ""),
                                summary_data=result.get("summary_data"),
                                source_info=src_info,
                                attachment_file_info=result.get("attachment_file_info", []),
                            )
                    if not NOTION_API_KEY:
                        st.caption("⚠️ NOTION_API_KEY 未設定のため保存不可")

            # マインドマップ
            tab_idx = 2
            if result.get("markmap_md"):
                with tabs[tab_idx]:
                    st.caption("💡 マウスホイールでズーム、ドラッグで移動できます。")
                    mm_html = render_markmap_html(result["markmap_md"])
                    st.components.v1.html(mm_html, height=560, scrolling=False)
                    st.download_button(
                        "📥 Markmap (.md)",
                        result["markmap_md"],
                        file_name=f"markmap_{fname_base}.md",
                        mime="text/markdown",
                        key=f"dmm_{result['date']}",
                    )
                tab_idx += 1

            # 構造化サマリー
            if result.get("summary_html") and len(tabs) > tab_idx:
                with tabs[tab_idx]:
                    st.info("💡 HTMLをダウンロードしてブラウザで開くと、見やすく印刷・PDF化できます。")
                    fname_base = (
                        Path(result["file_labels"][0]).stem if n_files == 1
                        else f"combined_{result['date'].replace(':','').replace(' ','_')}"
                    )
                    st.download_button(
                        "📥 構造化サマリー (.html)",
                        result["summary_html"],
                        file_name=f"summary_{fname_base}.html",
                        mime="text/html",
                        key=f"dsum_{result['date']}",
                    )
                    with st.expander("🔍 プレビュー（アプリ内）"):
                        st.components.v1.html(result["summary_html"], height=800, scrolling=True)

st.markdown("---")
st.caption("🎙️ 音声メモアプリ Pro ／ Powered by OpenAI Whisper & GPT-4o")
